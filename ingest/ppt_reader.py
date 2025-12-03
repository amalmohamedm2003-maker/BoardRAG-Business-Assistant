from pptx import Presentation
import os

def extract_pptx_text(path: str):
    prs = Presentation(path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        slides_data.append({
            "slide_number": i+1,
            "text": "\n".join(slide_text)
        })

    return {
        "file_name": os.path.basename(path),
        "slides": slides_data
    }
